import sys
import os 
import yaml
from PIL import Image, ImageOps
import tempfile
from pathlib import Path

sys.path.append('../../../DB')
import DB_utils

from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from datetime import datetime

import pandas as pd
import matplotlib
matplotlib.use("Agg")  # Tk 백엔드 비활성화(서버/스레드 환경 안전)
import matplotlib.pyplot as plt
import uuid
from copy import deepcopy

# image make all url 
def exchange_info_dict(d, *keys, default=""):
    for k in keys:
        if not isinstance(d, dict):
            return default
        d = d.get(k)
        if d is None:
            return default
    return d

def exchange_image_dict(d, *keys, default=""):
    for k in keys:
        # dict 접근
        if isinstance(d, dict):
            d = d.get(k, default)

        # list 접근
        elif isinstance(d, list) and isinstance(k, int):
            if 0 <= k < len(d):
                d = d[k]
            else:
                return default

        else:
            return default

        if d is None:
            return default

    return d



BASE_DIR = Path(__file__).resolve().parent
CONFIG_PATH = BASE_DIR / "configuration.yaml"
with open(CONFIG_PATH, "r", encoding="utf-8") as f:
    cfg = yaml.safe_load(f)


# 상단 이름 넣기 
def input_name_info(ppt,text : list):
    tf = ppt.text_frame 
    tf.clear()  # 기존 텍스트 제거
    lines = text
    
    for i, text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

            # ⭐ 문단 간 여백 제거
        run = p.add_run()
        run.text = text # 택스트 적용 
        run.font.name = "NanumGothic"
        
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0, 0, 0)
        # 줄별로 각각 어떤 포인트로 할 지 설정해야함
        if i == 0:    
            run.font.color.rgb = RGBColor(11, 50, 121)  # 파란색
        elif i== 4:
            run.font.bold = False
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.0    
            run.font.size = Pt(8)

        else:
            run.font.bold = False
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.0
                    
        # 정렬 
        p.alignment = PP_ALIGN.LEFT



def fix_image_orientation(img_path: str) -> Path:
    img = Image.open(img_path)
    img = ImageOps.exif_transpose(img)

    # ⭐ JPEG 대응 (RGBA → RGB)
    if img.mode in ("RGBA", "LA"):
        img = img.convert("RGB")

    tmp = tempfile.NamedTemporaryFile(
        suffix=Path(img_path).suffix,
        delete=False
    )
    img.save(tmp.name)
    tmp.close()

    return Path(tmp.name)
# 특정 공간에 이미지 넣기 공용 
def input_img(slide, ppt, img_path , width_rate = True , height_rate = True):
    fixed_img_path = fix_image_orientation(img_path)

    # 1. 기존 도형(네모칸)의 중심점 계산
    center_x = ppt.left + (ppt.width / 2)
    center_y = ppt.top + (ppt.height / 2)

    # 2. 일단 이미지를 삽입 (비율 유지를 위해 가로만 지정하거나 혹은 아예 지정 X)
    # 여기서는 기존 칸 너비에 맞추되 비율을 유지하도록 가로만 지정해봅니다.
    width_ = None
    height_ = None
    if width_rate :
        width_=ppt.width  # 가로 폭을 네모칸에 맞춤 (세로는 비율대로 자동 결정)
    if height_rate : 
        height_=ppt.height  # 가로 폭을 네모칸에 맞춤 (세로는 비율대로 자동 결정)

    new_img = slide.shapes.add_picture(
        str(fixed_img_path),
        ppt.left, 
        ppt.top,
        width=width_,  # 가로 폭을 네모칸에 맞춤 (세로는 비율대로 자동 결정)
        height=height_  # 가로 폭을 네모칸에 맞춤 (세로는 비율대로 자동 결정)
    )

    # 3. 삽입된 이미지의 새로운 중심점 기준으로 위치 재조정 (정중앙 정렬)
    new_img.left = int(center_x - (new_img.width / 2))
    new_img.top = int(center_y - (new_img.height / 2))

    # 4. (선택 사항) 만약 사진이 너무 길어서 네모칸 세로 범위를 벗어난다면?
    # 아래 로직을 추가하면 세로가 길 경우 세로 높이에 맞게 다시 줄여줍니다.
    if new_img.height > ppt.height:
        new_img.height = ppt.height
        new_img.width = int(new_img.width * (ppt.height / new_img.height)) # 비율 유지 축소
        # 크기가 바뀌었으니 다시 중앙 정렬
        new_img.left = int(center_x - (new_img.width / 2))
        new_img.top = int(center_y - (new_img.height / 2))

    # 5. 기존 도형 제거
    slide.shapes._spTree.remove(ppt._element)



# table 아래 테두리 
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.dml import MSO_THEME_COLOR


def set_all_cell_borders(cell, width=6350, color="BFBFBF" ,r = 0 ,c =0):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    setting_border_list = ['a:lnL', 'a:lnR','a:lnT', 'a:lnB']
    if c == 0 :
        setting_border_list.remove("a:lnL")
    else:
        setting_border_list.remove("a:lnR")

    for tag in setting_border_list:
        ln = OxmlElement(tag)
        ln.set('w', str(width))

        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', color)
        solidFill.append(srgbClr)

        prstDash = OxmlElement('a:prstDash')
        prstDash.set('val', 'solid')

        ln.append(solidFill)
        ln.append(prstDash)
        tcPr.append(ln)

def remove_side_border(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    lnL = OxmlElement('a:lnL')
    lnL.append(OxmlElement('a:noFill'))
    tcPr.append(lnL)

    lnR = OxmlElement('a:lnR')
    lnR.append(OxmlElement('a:noFill'))
    tcPr.append(lnR)

# table 관련 
def make_GI(slide , ppt, GI_data ,rows =13, cols = 2):
    left   = ppt.left
    top    = ppt.top
    width  = ppt.width
    height = ppt.height

    table = slide.shapes.add_table(
        rows, cols, left, top, width, height
    ).table

    table.first_row = False  # 머리글 행 특수 서식 해제
    table.first_col = False  # 첫 번째 열 특수 서식 해제
    # 표 너비를 도형 폭 기준 비율로 채움(좌:우 = 34:66)
    left_col_width = int(width * 0.36)
    table.columns[0].width = left_col_width
    table.columns[1].width = int(width - left_col_width)

    gray = RGBColor(231, 230, 230)  # 연한 회색 (배경)
    white = RGBColor(255, 255, 255) # 흰색 (배경)
    black = RGBColor(0, 0, 0)       # 검정색 (글자)

    cell_font_size = 8
    if rows > 40:
        cell_font_size = 4.5
    elif rows > 32:
        cell_font_size = 5
    elif rows > 26:
        cell_font_size = 6
    elif rows > 20:
        cell_font_size = 7

    for r_n, r in enumerate(range(rows)):
        # 모든 행의 높이를 일정하게 맞추고 싶다면 아래 주석 해제
        #table.rows[r].height = Cm(6.81) 
        #table.rows[r].width = Cm(8.04) 

        for c_n ,c in enumerate(range(cols)):
            cell = table.cell(r, c)
            set_all_cell_borders(table.cell(r, c) , r =r_n, c = c_n )
            cell.margin_top = Cm(0.0)
            cell.margin_bottom = Cm(0.0)
            # 배경색 설정
            cell.fill.solid()
            if c == 0:
                cell.fill.fore_color.rgb = gray
            else:
                cell.fill.fore_color.rgb = white

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # 폰트 설정 (폰트 8pt 적용)
            tf = cell.text_frame
            tf.clear()                      # ⭐ 중요
            paragraph = tf.paragraphs[0]
            if (c %2) == 0:
                paragraph.alignment = PP_ALIGN.CENTER  # 중앙 정렬 (필요시)
            else:
                paragraph.alignment = PP_ALIGN.LEFT
            # 텍스트가 이미 있거나 새로 넣을 때 폰트 적용을 위해 run 생성
            run = paragraph.add_run()
            run.font.name = "NanumGothic"
            run.font.size = Pt(cell_font_size)
            run.font.color.rgb = black
            
            # 예시 텍스트 입력 (내용이 필요할 경우)
            run.text = GI_data[r][c]

 
def make_GI_8cols(slide, ppt, GI_data):
    # GI_data: [[label, value], ...] -> 8열(키/값 x 4쌍) 테이블
    pair_per_row = 4
    rows = max(1, (len(GI_data) + pair_per_row - 1) // pair_per_row)
    cols = 8

    left = ppt.left
    top = ppt.top
    width = ppt.width
    height = ppt.height

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.first_row = False
    table.first_col = False

    # 열 비율: (키 8% + 값 17%) x 4 = 100%
    key_w = int(width * 0.08)
    val_w = int(width * 0.17)
    used = 0
    for i in range(cols):
        if i == cols - 1:
            table.columns[i].width = int(width - used)
        else:
            col_w = key_w if i % 2 == 0 else val_w
            table.columns[i].width = col_w
            used += col_w

    gray = RGBColor(231, 230, 230)
    white = RGBColor(255, 255, 255)
    black = RGBColor(0, 0, 0)

    cell_font_size = 9
    if rows > 18:
        cell_font_size = 6
    elif rows > 14:
        cell_font_size = 7
    elif rows > 10:
        cell_font_size = 8

    for r in range(rows):
        base_idx = r * pair_per_row
        row_values = ["-"] * cols
        for pair_idx in range(pair_per_row):
            src_idx = base_idx + pair_idx
            if src_idx >= len(GI_data):
                break
            row_values[pair_idx * 2] = GI_data[src_idx][0]
            row_values[pair_idx * 2 + 1] = GI_data[src_idx][1]

        for c in range(cols):
            cell = table.cell(r, c)
            set_all_cell_borders(cell, r=r, c=c)
            cell.margin_top = Cm(0.0)
            cell.margin_bottom = Cm(0.0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = gray if c % 2 == 0 else white
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            tf = cell.text_frame
            tf.clear()
            paragraph = tf.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER if c % 2 == 0 else PP_ALIGN.LEFT

            run = paragraph.add_run()
            run.font.name = "NanumGothic"
            run.font.size = Pt(cell_font_size)
            run.font.color.rgb = black
            run.text = str(row_values[c])

# def Features 
def featrue_info(ppt, text):
    ppt.width = Cm(8.16)
    tf = ppt.text_frame
    tf.clear()  # 기존 서식 제거

    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = text  # 입력받은 텍스트 리스트

    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # ⭐ 불렛 포인트(동그라미) 설정
        p.level = 0  # 불렛 수준 설정 (기본 0)
        
        # 문단 간 여백 및 줄 간격 제거 (촘촘하게)
        #p.space_before = Pt(1)
        p.line_spacing = 1.3
        
        # 정렬: 왼쪽 정렬
        p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        # 텍스트 앞에 동그라미 기호 추가 (가장 확실한 방법)
        run.text = f"• {line_text}" 
        
        # 폰트 기본 설정 (나눔고딕, 10pt)
        run.font.name = "NanumGothic"
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = RGBColor(0, 0, 0) # 기본 검정색

# 텍스트 넣기 

color_define = {
    "red" : (205,35,50),
    "blue" : (11,50,121),
    "white" : (255,255,255)
}
def input_oneline_text(ppt,text,color = "red", fsize = 28 , bold = False ):
    tf = ppt.text_frame 
    tf.clear()  # 기존 텍스트 제거
    lines = text
    
    for i, text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            # 중 추가 
            p = tf.add_paragraph()

            # ⭐ 문단 간 여백 제거
        run = p.add_run()
        run.text = text # 택스트 적용 
        run.font.name = "NanumGothic"
        run.font.size = Pt(fsize)

        r, g, b = color_define[color]
        run.font.color.rgb = RGBColor(r,g,b)
        # 줄별로 각각 어떤 포인트로 할 지 설정해야함
       
        run.font.bold = bold              
        # 정렬 
        p.alignment = PP_ALIGN.RIGHT

def input_title(ppt, text):
    tf = ppt.text_frame
    tf.clear()                      # 기존 서식 제거

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text

    run.font.name = "NanumGothic"   # 나눔고딕으로 변경 
    run.font.size = Pt(28)          # font 변경 
    run.font.bold = True            # ✅ 볼드



# =========================
# 2. 합계 계산 및 포맷팅
# =========================
def calculate_total(df , sum_targets):
    total_row = {col: "" for col in df.columns}
    total_row["층"] = "합계"
    for col in sum_targets:
        if col in df.columns:
            # 콤마 제거 후 수치화하여 합산
            vals = pd.to_numeric(df[col].astype(str).str.replace(',', '').str.strip(), errors='coerce')
            total_row[col] = vals.sum()
    return pd.DataFrame([total_row])

def format_with_commas(val):
    try:
        if isinstance(val, (int, float)) and not pd.isna(val) and val != "":
            return f"{val:,.2f}".rstrip('0').rstrip('.') if val % 1 != 0 else f"{val:,.0f}"
        return val
    except:
        return val

def get_optimal_widths(df):
    widths = []
    for col in df.columns:
        # 데이터와 컬럼명 중 더 긴 것 기준으로 기본 너비 계산
        max_str_len = max(df[col].astype(str).map(len).max(), len(col))
        
        if col == "비고":
            # 비고 칸은 최소 10글자(한글/영문 혼합 고려 약 2.5~3.0인치) 이상 확보
            w = max(max_str_len * 0.2, 3.0) 
        else:
            # 다른 일반 컬럼들의 기본 너비
            w = max(max_str_len * 0.2, 1.0)
            
        widths.append(w)
    return widths
    

######### detail 표 생성 
def make_detail_table(datas , save_path):
    # =========================
    # 1. 데이터
    # =========================
    data = datas
    df = pd.DataFrame(data)
    if df.shape[1] > 9:
        df = df.iloc[:, :9]

    # 사용자가 요청한 순서로 컬럼명 매핑 및 정렬
    # (기존 데이터의 순서와 상관없이 아래 순서로 재배치됩니다)
    ordered_columns = ["층", "업종", "면적(평)", "면적(㎡)", "보증금", "월임대료", "관리비", "비고", "공실여부"]
    if df.shape[1] >= 9:
        df.columns = ["층", "업종", "보증금", "관리비", "면적(평)", "면적(㎡)", "월임대료", "비고", "공실여부"]
    else:
        df.columns = ["층", "업종", "보증금", "관리비", "면적(평)", "면적(㎡)", "월임대료", "비고"]
        df["공실여부"] = "-"
    df = df[ordered_columns] # 순서 재정렬

    # 합계 대상 컬럼 (숫자 데이터)
    sum_targets = ["면적(평)", "면적(㎡)", "보증금", "월임대료", "관리비"]


    df_with_total = pd.concat([df, calculate_total(df , sum_targets)], ignore_index=True)
    for col in sum_targets:
        df_with_total[col] = df_with_total[col].apply(format_with_commas)

    df_with_total = df_with_total.fillna("-").replace("", "-")

    # =========================
    # 3. 스타일 및 여백 제거 설정
    # =========================
    plt.rcParams["font.family"] = "Malgun Gothic"
    plt.rcParams["axes.unicode_minus"] = False



    col_widths_inches = get_optimal_widths(df_with_total)
    total_width = sum(col_widths_inches)
    fig_height = (len(df_with_total) + 1) * 0.45

    fig, ax = plt.subplots(figsize=(total_width, fig_height))
    ax.axis("off")

    # 테이블 생성
    table = ax.table(
        cellText=df_with_total.values,
        colLabels=df_with_total.columns,
        colWidths=[w/total_width for w in col_widths_inches],
        cellLoc="center",
        loc="center",
        bbox=[0, 0, 1, 1]
    )

    table.auto_set_font_size(False)

    # 셀 스타일링
    last_row_idx = len(df_with_total)
    for (row, col), cell in table.get_celld().items():
        if row == 0: # 헤더
            cell.set_text_props(weight="bold", size=12, color="#FFFFFF")
            cell.set_facecolor("#374151")
        elif row == last_row_idx: # 합계 행
            cell.set_text_props(weight="bold", size=11, color="#1D4ED8")
            cell.set_facecolor("#EFF6FF")
        else: # 일반 행
            cell.set_text_props(size=11, color="#111827")
            cell.set_facecolor("#FFFFFF" if row % 2 else "#F9FAFB")
        cell.set_edgecolor("#D1D5DB")



    filename = f"lease_table_{uuid.uuid4().hex}.png"

    plt.savefig(
        save_path / filename,
        dpi=220,
        bbox_inches="tight",
        pad_inches=0
    )
    plt.close()
    return save_path / filename


# 슬라이드 복사 
def duplicate_slide(prs, slide):
    """
    slide 하나를 그대로 복사해서 prs에 추가
    """
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    for shape in slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(
            new_el, 'p:extLst'
        )
    return new_slide

def move_slide(prs, from_idx, to_idx):
    slides = prs.slides
    slide_ids = slides._sldIdLst

    slide_id = slide_ids[from_idx]
    slide_ids.remove(slide_id)
    slide_ids.insert(to_idx, slide_id)

# 마지막 슬라이드 제거 
def remove_last_slide(prs , number = -1 ):
    slides = prs.slides
    slide_id_list = slides._sldIdLst  # 내부 XML
    slide_id_list.remove(slide_id_list[number])

def safe(value):
    return str(value) if value else "- "

LEFT_BOOLEAN_FIELDS = {
    "is_new_site",
    "is_remodeling",
    "is_office_building",
    "is_investment",
    "is_development",
    "is_stable_holding",
}

LEFT_SIDEBAR_FIELDS = [
    ("입지판단", "location_decide"),
    ("가격판단", "price_decide"),
    ("수익률판단", "yield_decide"),
    ("명도판단", "vacancy_decide"),
    ("제한판단", "limit_decide"),
    ("상태판단", "loan_decide"),
    ("건물명", "bd_name"),
    ("주소", "address"),
    ("상세주소", "address_detail"),
    ("주변역", "nearby_station"),
    ("주변역2", "nearby_station2"),
    ("입지", "site_location"),
    ("입금가(만원)", "deposit_price"),
    ("기타매매가(만원)", "other_sale_price"),
    ("매매가(만원)", "sale_price"),
    ("수익률(%)", "yield_rate"),
    ("보증금(만원)", "security_deposit"),
    ("토지평단가(만원)", "price_per_pyeong"),
    ("월임대료(만원)", "monthly_rent_fee"),
    ("연면적평단가(만원)", "price_per_total_floor_area"),
    ("관리비(만원)", "maintenance_fee"),
    ("관리비지출(만원)", "maintenance_expense"),
    ("대출현황", "loan_status"),
    ("토지지목", "land_category"),
    ("용도지역", "zoning_type"),
    ("토지면적(평)", "land_area_pyeong"),
    ("토지면적(㎡)", "land_area_sqm"),
    ("㎡당공시(원)", "official_price_per_sqm_won"),
    ("공시연도", "official_price_per_pyeong_million_date"),
    ("공시지가평당(만원)", "official_price_per_pyeong_million"),
    ("공시합계(만원)", "official_price_total_million"),
    ("도로", "road_access"),
    ("도로2", "road_access2"),
    ("승인일자", "approval_date"),
    ("냉난방", "heating_cooling"),
    ("위반사항", "violation_info"),
    ("건축용도", "building_usage"),
    ("건물구조", "building_structure"),
    ("연면적(평)", "gross_area_pyeong"),
    ("연면적(㎡)", "gross_area_sqm"),
    ("사용가능면적(평)", "usable_area_pyeong"),
    ("건축면적(평)", "building_area_pyeong"),
    ("건축면적(㎡)", "building_area_sqm"),
    ("층고(m)", "floor_height"),
    ("보높이(m)", "beam_clearance_height"),
    ("지하층", "underground_floors"),
    ("지상층", "aboveground_floors"),
    ("승강기", "elevator"),
    ("비상승강기", "emergency_elevator"),
    ("건폐율(%)", "building_coverage_ratio"),
    ("용적률(%)", "floor_area_ratio"),
    ("총주차", "parking_capacity"),
    ("옥외기계", "parking_outdoor_mechanical"),
    ("옥외자주", "parking_outdoor_self"),
    ("옥내기계", "parking_indoor_mechanical"),
    ("옥내자주", "parking_indoor_self"),
    ("방향기준", "direction_basis"),
    ("방향", "direction"),
    ("방개수", "room_count"),
    ("화장실개수", "bathroom_count"),
    ("신축부지", "is_new_site"),
    ("리모델링", "is_remodeling"),
    ("사옥형", "is_office_building"),
    ("수익형", "is_investment"),
    ("개발/전환", "is_development"),
    ("보유안정", "is_stable_holding"),
]

def _format_left_sidebar_value(field_key, value):
    if field_key in LEFT_BOOLEAN_FIELDS:
        normalized = str(value).strip().lower()
        return "Y" if normalized in {"true", "1", "y", "yes"} else "-"

    if value is None:
        return "-"

    text = str(value).strip()
    if text == "":
        return "-"

    if field_key == "approval_date":
        return text.split(" ")[0]

    return text

def build_left_sidebar_gi_data(building_info):
    if not isinstance(building_info, dict):
        return [["좌측정보", "-"]]

    rows = []
    for label, field_key in LEFT_SIDEBAR_FIELDS:
        rows.append([label, _format_left_sidebar_value(field_key, building_info.get(field_key))])
    return rows

import re 
# 
def sanitize_filename(name: str) -> str:
    name = name.strip()                # 1️⃣ 양 끝 정리
    name = name.replace('\t', ' ')     # 2️⃣ 탭 제거
    name = name.replace('\n', ' ')     # 3️⃣ 줄바꿈 제거
    # name = re.sub(r'[\\/:*?"<>|]', '', name)  # 4️⃣ 금지 문자 제거
    # name = re.sub(r'\s+', ' ', name)   # 5️⃣ 공백 정리
    return name

def run(bd_number):

    ### 해당 정보를 바탕으로 ppt 만들기 진행 
    conn = DB_utils.join_db()
    results = DB_utils.ppt_info_search(conn, bd_number)
    # 페이지 3,4 이상 넣어야할 이미지 
    p34_images = results["images"]
    print("시작~~~~")
    ## 필요 값 정의
    ## 모두 값이 없을 수 있음

    filename_askprice = "-억원"
    bd_address = exchange_info_dict(results, "building_info", "address")
    bd_name = exchange_info_dict(results, "building_info", "bd_name")

    middle_img1 = exchange_image_dict(results, "images", "main", 0, "url")
    middle_img2 = exchange_image_dict(results, "images", "sub1", 0, "url")
    middle_img3 = exchange_image_dict(results, "images", "sub2", 0, "url")

    gi1 = exchange_info_dict(results, "building_info", "address")
    gi1_1 = exchange_info_dict(results, "building_info", "address_detail")
    gi2 = exchange_info_dict(results, "building_info", "zoning_type")
    gi3 = exchange_info_dict(results, "building_info", "land_area_sqm")
    gi3_1 =exchange_info_dict(results, "building_info", "land_area_pyeong")

    gi4 = exchange_info_dict(results, "building_info", "gross_area_sqm")
    gi4_1 = exchange_info_dict(results, "building_info", "gross_area_pyeong")

    gi5 = exchange_info_dict(results, "building_info", "building_coverage_ratio")
    gi6 = exchange_info_dict(results, "building_info", "floor_area_ratio")

    gi7 = exchange_info_dict(results, "building_info", "aboveground_floors")
    gi8 = exchange_info_dict(results, "building_info", "underground_floors")

    gi9 = exchange_info_dict(results, "building_info", "building_usage")
    gi10 = exchange_info_dict(results, "building_info", "building_structure")

    gi11 = exchange_info_dict(results, "building_info", "parking_capacity")
    gi11_1 = safe(exchange_info_dict(results, "building_info", "parking_indoor_mechanical"))
    gi11_2 = safe(exchange_info_dict(results, "building_info", "parking_indoor_self"))
    gi11_3 = safe(exchange_info_dict(results, "building_info", "parking_outdoor_mechanical"))
    gi11_4 = safe(exchange_info_dict(results, "building_info", "parking_outdoor_self"))
    
    
    
    gi12 = exchange_info_dict(results, "building_info", "elevator")
    gi13 = exchange_info_dict(results, "building_info", "emergency_elevator")

    gi14 = exchange_info_dict(results, "building_info", "direction")
    gi15 = exchange_info_dict(results, "building_info", "approval_date")
    gi16 = exchange_info_dict(results, "building_info", "official_price_per_pyeong_million")
    gi17= exchange_info_dict(results, "building_info", "official_price_per_pyeong_million_date")

    feature_text = exchange_info_dict(results, "building_memo", "bd_feature").split("\n")

    sale_p = safe(exchange_info_dict(results, "building_info", "sale_price"))
    right1 = safe(exchange_info_dict(results, "building_info", "price_per_pyeong"))
    right2 = safe(exchange_info_dict(results, "building_info", "price_per_total_floor_area"))

    right3 = safe(exchange_info_dict(results, "building_info", "security_deposit"))
    right4 = safe(exchange_info_dict(results, "building_info", "monthly_rent_fee"))
    right5 = safe(exchange_info_dict(results, "building_info", "maintenance_fee"))
    right6 = safe(exchange_info_dict(results, "building_info", "yield_rate"))

    page4_img1 = exchange_image_dict(results, "images", "sub3", 0, "url")
    page4_img2 = exchange_image_dict(results, "images", "sub4", 0, "url")
    
    filji_text =""
    filji_count = len(gi1_1.split("/"))

    if (safe(gi1_1) != "- "):
        filji_text=f"\n 외 {filji_count}필지"
    
    if (safe(gi12) == "- ") and (safe(gi13) == "- "):
        gi12_13 = "-"
    elif (safe(gi12) != "- ") and (safe(gi13) == "- "):
        gi12_13 = f"승용 {gi12}대"
    elif (safe(gi12) == "- ") and (safe(gi13) != "- "):
        gi12_13 = f"비상 {gi13}대"
    else:
        gi12_13 =f"승용 {safe(gi12)}대 / 비상 {safe(gi13)}대"

    new_gi15 = safe(gi15).split(".")
    if len(new_gi15) > 1 :
        result_gi15 = f"{new_gi15[0]}년 {new_gi15[1]}월 {new_gi15[2]}일"
    else:
        result_gi15 = safe(gi15)

    if safe(gi11) in ["0", "- "]:
        result_gi11 = "-"
    else:
        result_gi11 = f"{safe(gi11)}대("

        if gi11_1 != "- ":
            result_gi11 += f"옥내 기계 {gi11_1}"
        if gi11_2 != "- ":
            result_gi11 += f", 옥내 자주 {gi11_2}"
        if gi11_3 != "- ":
            result_gi11 += f", 옥외 기계 {gi11_3}"
        if gi11_4 != "- ":
            result_gi11 += f", 옥외 자주 {gi11_4}"
        
        result_gi11 += ")"

        result_gi11 = result_gi11.replace("(, ","(")
        
    GI_data = [
        ["대지위치",safe(gi1)+"번지" + filji_text],
        ["지역지구",safe(gi2)],
        ["대지면적",f"{safe(gi3)}㎡ ({safe(gi3_1)}평)"],
        ["연 면 적",f"{safe(gi4)}㎡ ({safe(gi4_1)}평)"],
        ["건폐율/용적률", f"{safe(gi5)}% / {safe(gi6)}%"],
        ["건물규모",f"지상{safe(gi7)}층 / 지하 {safe(gi8)}층"],
        ["건축물주용도",safe(gi9)],
        ["건축물주구조",safe(gi10)],
        ["주차대수",f"{result_gi11}"],
        ["승 강 기", gi12_13],
        ["방향(주출입구기준)",f"{safe(gi14)}"],
        ["사용승인일",result_gi15],
        ["공시지가",f"{safe(gi16)}만원 / 3.3㎡ ({safe(gi17)}년 기준)"]
    ]

    # ppt load 
    ppt_BASE_DIR = Path(__file__).resolve().parent.parent  # ppt/
    ppt_path = ppt_BASE_DIR / "statics" / "template.pptx"
    ppt_save_path = ppt_BASE_DIR / "statics"
    prs = Presentation(ppt_path)

    # 새 첫 페이지 생성: 기존 첫 슬라이드를 복제해 맨 앞으로 이동
    left_sidebar_gi_data = build_left_sidebar_gi_data(results.get("building_info") or {})
    new_first_slide = duplicate_slide(prs, prs.slides[0])
    move_slide(prs, len(prs.slides) - 1, 0)

    # 새 1페이지는 중앙/부가 요소를 모두 제거하고 표만 크게 채움
    gi_shape = None
    for shape in list(new_first_slide.shapes):
        if shape.name == "GI_table":
            gi_shape = shape
            continue
        new_first_slide.shapes._spTree.remove(shape._element)

    if gi_shape is not None:
        gi_shape.left = Cm(0.5)
        gi_shape.top = Cm(0.5)
        gi_shape.width = prs.slide_width - Cm(1.0)
        gi_shape.height = prs.slide_height - Cm(1.0)
        make_GI_8cols(new_first_slide, gi_shape, left_sidebar_gi_data)

    # 대상 슬라이드 설정 
    slide = prs.slides[1] # 대상 슬라이드 찾기 

    # 지정된 이름의 객체 선택 
    for shape in slide.shapes:
        print("shape:",shape)
        
        # title 설정 
        if shape.name == "title":
            input_title(shape,bd_name)
            
        ## main / sub person name 넣기
        elif shape.name == "main_person":
            input_name_info(shape,cfg["main_person"])
        elif shape.name == "sub_person":
            input_name_info(shape, cfg["sub_person"])
        
        ## middle image 입력 
        elif shape.name == "middle_img1":
            if middle_img1 != "":
                input_img(slide,shape, cfg["path"] + middle_img1)
        elif shape.name == "middle_img2":
            if middle_img2 != "":
                input_img(slide,shape, cfg["path"] + middle_img2)
        elif shape.name == "middle_img3":
            if middle_img3 != "":
                input_img(slide,shape, cfg["path"] + middle_img3)
            
            
            
        # GI 입력
        elif shape.name == "GI_table":
            print("dddd")
            make_GI(slide, shape , GI_data)
            
        elif shape.name == "features":
            featrue_info(shape, feature_text)


        elif shape.name =="ask_price":
            price = int(sale_p.replace(",", ""))
            value = price / 10000

            # 소수 첫째자리가 0이면 제거
            if value.is_integer():
                text = f"{int(value):,} 억원"
            else:
                text = f"{value:,.1f} 억원"
            input_oneline_text(shape, [f"{text}"], color = 'red', fsize=28 , bold = True)

            filename_askprice = text

        elif shape.name =="area_price":
            input_oneline_text(shape, [f"{right1} 만원/3.3㎡",f"{right2} 만원/3.3㎡"],color = 'blue', fsize= 10 )
        elif shape.name =="rent_price":
            input_oneline_text(shape, [f"{right3} 만원",f"{right4} 만원",f"{right5} 만원", f"{right6}%"],color = 'blue', fsize= 10)

        elif shape.name =="date":

            now = datetime.now()
            date_str = now.strftime("%Y년 %m월")
            input_oneline_text(shape, [date_str],color = 'white', fsize= 9 , bold= True)


    ##################### page 2  
    slide_number = 2 
    slide2_count = (len(p34_images.get("sub5", [])) + 5) // 6

    for now_slide_number in range(slide2_count):

        slide = prs.slides[slide_number] # 대상 슬라이드 찾기 
        if now_slide_number != (slide2_count -1 ):
            new_slide = duplicate_slide(prs, prs.slides[slide_number])
            move_slide(prs, len(prs.slides) - 1, slide_number+ 1 )
        page2_img1 = exchange_image_dict(results, "images", "sub5", (6 *(now_slide_number + 1) ) - 6, "url")
        page2_img2 = exchange_image_dict(results, "images", "sub5", (6 *(now_slide_number + 1) ) - 5, "url")
        page2_img3 = exchange_image_dict(results, "images", "sub5", (6 *(now_slide_number + 1) ) - 4, "url")
        page2_img4 = exchange_image_dict(results, "images", "sub5", (6 *(now_slide_number + 1) ) - 3, "url")
        page2_img5 = exchange_image_dict(results, "images", "sub5", (6 *(now_slide_number + 1) ) - 2, "url")
        page2_img6 = exchange_image_dict(results, "images", "sub5", (6 *(now_slide_number + 1) ) - 1, "url")

        # 지정된 이름의 객체 선택 
        for shape in slide.shapes:
            print("shape:",shape)
            
            ## main / sub person name 넣기
            if shape.name == "main_person":
                input_name_info(shape, cfg["main_person"])
            elif shape.name == "sub_person":
                input_name_info(shape, cfg["sub_person"])
            
            ## middle image 입력 
            elif (shape.name == "img1") and (page2_img1 != ""):
                input_img(slide,shape, cfg["path"] + page2_img1)
            elif (shape.name == "img2") and (page2_img2 != ""):
                input_img(slide,shape, cfg["path"] + page2_img2)
            elif (shape.name == "img3") and (page2_img3 != ""):
                input_img(slide,shape, cfg["path"] + page2_img3)
            elif (shape.name == "img4") and (page2_img4 != ""):
                input_img(slide,shape, cfg["path"] + page2_img4)
            elif (shape.name == "img5") and (page2_img5 != ""):
                input_img(slide,shape, cfg["path"] + page2_img5)
            elif (shape.name == "img6") and (page2_img6 != ""):
                input_img(slide,shape, cfg["path"] + page2_img6)
        
        slide_number += 1
    if len(p34_images.get("sub5", [])) == 0:
        slide_number += 1
    # 마지막 하나 생성되었을거니 지운다 
    # remove_last_slide(prs , slide_number  )
    ##################### page 3  
    slide = prs.slides[slide_number] # 대상 슬라이드 찾기 
    results = DB_utils.extract_detail_management(conn,bd_number)
    table_img_path = make_detail_table(results , ppt_save_path)
    # 지정된 이름의 객체 선택 
    for shape in slide.shapes:
        print("shape:",shape)
        
        ## main / sub person name 넣기
        if shape.name == "main_person":
            input_name_info(shape, cfg["main_person"])
        elif shape.name == "sub_person":
            input_name_info(shape, cfg["sub_person"])
        
        # middle image 입력 
        elif shape.name == "img1":
            print("table_img_path" , table_img_path)
            input_img(slide,shape, f"{table_img_path}" , True, False)

    slide_number += 1 
    ##################### page 4 
    save_img_list = []
    save_title_list = []

    for i in p34_images.get("sub3", []):
        save_img_list.append(i["url"])
        save_title_list.append("토지이용계획")

    for i in p34_images.get("sub4", []):
        save_img_list.append(i["url"])
        save_title_list.append("건축물대장")

    for i in p34_images.get("sub6", []):
        save_img_list.append(i["url"])
        save_title_list.append("참고자료")
    print("slide_number" , slide_number)
    for i in range(len(save_img_list)):
        slide = prs.slides[slide_number + i] # 대상 슬라이드 찾기 

        # 슬라이드 추가 
        duplicate_slide(prs, slide)

        # 지정된 이름의 객체 선택 
        for shape in slide.shapes:
            print("shape:",shape)
            
            ## main / sub person name 넣기
            if shape.name == "main_person":
                input_name_info(shape, cfg["main_person"])
            elif shape.name == "sub_person":
                input_name_info(shape, cfg["sub_person"])
            elif shape.name == "title":
                input_title(shape,save_title_list[i])
            
            ## middle image 입력 
            elif (shape.name == "img_1"):
                input_img(slide,shape, cfg["path"] + save_img_list[i])

    # 마지막 하나는 제거한다.
    remove_last_slide(prs)
    # 원하는 형식의 이름으로 변경
    rename_list = bd_address.strip().split(" ")
    rename_address = rename_list[-2] + " "+rename_list[-1]
    filename = sanitize_filename(f"[ERA]매매{filename_askprice.replace(' ','')}_{rename_address}_{bd_name}(유지혜D).pptx")
    prs.save(ppt_save_path/ filename)

    return ppt_save_path/ filename , filename
