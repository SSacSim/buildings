import re
import sys
import uuid
from io import BytesIO
from copy import deepcopy
from datetime import datetime
from pathlib import Path

import yaml
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

sys.path.append("../../../DB")
import DB_utils  # noqa: E402


BASE_DIR = Path(__file__).resolve().parent
PPT_DIR = BASE_DIR.parent
CONFIG_PATH = BASE_DIR / "configuration.yaml"

with open(CONFIG_PATH, "r", encoding="utf-8") as f:
    cfg = yaml.safe_load(f)


ROW_KEYS = [
    "address",
    "zoning_type",
    "land_area",
    "gross_area",
    "building_area",
    "floors",
    "parking_elevator",
    "approval_date",
    "official_price",
    "sale_price",
    "price_per_pyeong",
    "price_per_total",
    "usable_area",
    "yield_rate",
]


def _safe(value):
    if value is None:
        return "-"
    text = str(value).strip()
    return text if text else "-"


def _to_number(value):
    text = _safe(value).replace(",", "")
    if text == "-":
        return None
    try:
        return float(text)
    except Exception:
        return None


def _fmt_with_unit(value, unit):
    text = _safe(value)
    if text == "-":
        return "-"
    return f"{text}{unit}"


def _fmt_approval_date(value):
    text = _safe(value)
    if text == "-":
        return "-"
    return text.split(" ")[0]


def _fmt_sale_price_eok(value):
    num = _to_number(value)
    if num is None:
        return "-"
    eok = num / 10000.0
    if eok.is_integer():
        return f"{int(eok):,}억"
    return f"{eok:,.1f}억"


def _fmt_price_per_area(value):
    text = _safe(value)
    if text == "-":
        return "-"
    return f"{text}만원 / 3.3㎡"


def _fmt_count(value):
    text = _safe(value)
    if text == "-":
        return "-"
    digits = re.findall(r"\d+", text)
    if digits:
        return f"{digits[0]}대"
    return text if text.endswith("대") else f"{text}대"


def _build_rows(building_info):
    above = _safe(building_info.get("aboveground_floors"))
    below = _safe(building_info.get("underground_floors"))
    if above == "-" and below == "-":
        floors = "-"
    else:
        floors = f"지상 {above}층 / 지하 {below}층"

    parking = _fmt_count(building_info.get("parking_capacity"))
    elevator = _fmt_count(building_info.get("elevator"))
    parking_elevator = "-" if parking == "-" and elevator == "-" else f"{parking} / {elevator}"

    data = {
        "address": _safe(building_info.get("address")),
        "zoning_type": _safe(building_info.get("zoning_type")),
        "land_area": _fmt_with_unit(building_info.get("land_area_pyeong"), "평"),
        "gross_area": _fmt_with_unit(building_info.get("gross_area_pyeong"), "평"),
        "building_area": _fmt_with_unit(building_info.get("building_area_pyeong"), "평"),
        "floors": floors,
        "parking_elevator": parking_elevator,
        "approval_date": _fmt_approval_date(building_info.get("approval_date")),
        "official_price": _fmt_price_per_area(building_info.get("official_price_per_pyeong_million")),
        "sale_price": _fmt_sale_price_eok(building_info.get("sale_price")),
        "price_per_pyeong": _fmt_price_per_area(building_info.get("price_per_pyeong")),
        "price_per_total": _fmt_price_per_area(building_info.get("price_per_total_floor_area")),
        "usable_area": _fmt_with_unit(building_info.get("usable_area_pyeong"), "평"),
        "yield_rate": _fmt_with_unit(building_info.get("yield_rate"), "%"),
    }
    return [data[k] for k in ROW_KEYS]


def _iter_shapes(shapes):
    for shape in shapes:
        if hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)
        yield shape


def _find_shape(slide, name):
    target = name.strip().lower()
    for shape in _iter_shapes(slide.shapes):
        if (shape.name or "").strip().lower() == target:
            return shape
    return None


def _set_shape_text(shape, text, size_pt, bold=False):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(text)
    for run in p.runs:
        run.font.name = "NanumGothic"
        run.font.size = Pt(size_pt)
        run.font.bold = bold


def _set_cell_text(cell, text, size_pt):
    cell.text = str(text)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = 0
    cell.margin_bottom = 0
    tf = cell.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = "NanumGothic"
            run.font.size = Pt(size_pt)


def _fix_image_orientation(img_path):
    img = Image.open(img_path)
    img = ImageOps.exif_transpose(img)
    if img.mode in ("RGBA", "LA"):
        img = img.convert("RGB")
    tmp_path = BASE_DIR / f"_tmp_{uuid.uuid4().hex}{Path(img_path).suffix}"
    img.save(tmp_path)
    return tmp_path


def _replace_picture(slide, placeholder_shape, image_path):
    fixed_path = _fix_image_orientation(image_path)
    try:
        new_pic = slide.shapes.add_picture(
            str(fixed_path),
            placeholder_shape.left,
            placeholder_shape.top,
            width=placeholder_shape.width,
            height=placeholder_shape.height,
        )
        sp_tree = slide.shapes._spTree
        sp_tree.remove(new_pic._element)
        sp_tree.insert(sp_tree.index(placeholder_shape._element), new_pic._element)
        sp_tree.remove(placeholder_shape._element)
    finally:
        try:
            fixed_path.unlink(missing_ok=True)
        except Exception:
            pass


def _to_local_image_path(url):
    backend_root = Path(cfg.get("path", "")).resolve()
    return str((backend_root / url.lstrip("/")).resolve())


def _duplicate_slide(prs, source_slide):
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for shape in source_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                new_slide.shapes.add_picture(
                    BytesIO(shape.image.blob),
                    shape.left,
                    shape.top,
                    width=shape.width,
                    height=shape.height,
                )
                continue
            except Exception:
                pass

        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return new_slide


def _fill_slide(slide, infos):
    # n1~n7
    for i in range(1, 8):
        shape = _find_shape(slide, f"n{i}")
        if not shape:
            continue
        if i <= len(infos):
            bd_name = _safe((infos[i - 1].get("building_info") or {}).get("bd_name"))
            _set_shape_text(shape, bd_name, 11, True)
        else:
            _set_shape_text(shape, "-", 11, True)

    # p1~p7
    for i in range(1, 8):
        shape = _find_shape(slide, f"p{i}")
        if not shape or i > len(infos):
            continue
        main_images = (infos[i - 1].get("images") or {}).get("main", [])
        if not main_images:
            continue
        image_url = main_images[0].get("url")
        if not image_url:
            continue
        image_path = _to_local_image_path(image_url)
        if Path(image_path).exists():
            _replace_picture(slide, shape, image_path)

    # table
    table_shape = _find_shape(slide, "table")
    if table_shape is None or not getattr(table_shape, "has_table", False):
        return

    table = table_shape.table
    max_rows = min(len(ROW_KEYS), len(table.rows))
    max_cols = max((len(r.cells) for r in table.rows), default=0)
    start_col = 1 if max_cols >= 8 else 0
    fill_cols = min(7, max(0, max_cols - start_col))

    for offset in range(fill_cols):
        col_idx = start_col + offset
        values = []
        if offset < len(infos):
            building_info = infos[offset].get("building_info") or {}
            values = _build_rows(building_info)
        for row_idx in range(max_rows):
            row_cells = table.rows[row_idx].cells
            if col_idx >= len(row_cells):
                continue
            value = values[row_idx] if row_idx < len(values) else "-"
            _set_cell_text(row_cells[col_idx], value, 8)


def run(bd_numbers):
    if not bd_numbers:
        raise ValueError("bd_numbers is empty")

    normalized_ids = [int(x) for x in bd_numbers if x is not None]
    if not normalized_ids:
        raise ValueError("bd_numbers is empty")

    conn = DB_utils.join_db()
    try:
        infos_all = [DB_utils.ppt_info_search(conn, bd) for bd in normalized_ids]
    finally:
        conn.close()

    template_path = PPT_DIR / "statics" / "template2.pptx"
    prs = Presentation(str(template_path))
    template_slide = prs.slides[0]

    chunks = [infos_all[i : i + 7] for i in range(0, len(infos_all), 7)]
    if not chunks:
        chunks = [[]]

    slides_to_fill = [template_slide]
    for _ in chunks[1:]:
        slides_to_fill.append(_duplicate_slide(prs, template_slide))

    for idx, chunk in enumerate(chunks):
        _fill_slide(slides_to_fill[idx], chunk)

    out_name = f"compare_{datetime.now().strftime('%Y%m%d')}_{uuid.uuid4().hex[:6]}.pptx"
    out_path = PPT_DIR / "statics" / out_name
    prs.save(str(out_path))
    return str(out_path), out_name


if __name__ == "__main__":
    args = [int(x) for x in sys.argv[1:]]
    ppt_path, file_name = run(args)
    print(ppt_path)
    print(file_name)
